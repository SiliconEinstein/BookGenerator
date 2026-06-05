from __future__ import annotations
import json
from pathlib import Path
from pptx import Presentation
from app.services.llm_service import gpt_completion, parse_json_object
from app.utils.fs_utils import ensure_dir
from app.utils.prompt_utils import render_prompt


def _language_label(output_language: int) -> str:
    return '英文' if output_language == 2 else '中文'


async def plan_ppt_slides(
    content: dict,
    chapter_name: str,
    section_name: str,
    output_language: int = 1,
) -> dict:
    prompt = render_prompt(
        'ppt_plan',
        chapter_name=chapter_name,
        section_name=section_name,
        section_content=json.dumps(content, ensure_ascii=False, indent=2),
        output_language_label=_language_label(output_language),
    )
    plan = parse_json_object(await gpt_completion(prompt, temperature=0.5, max_tokens=6000))
    return _normalize_ppt_plan(plan, chapter_name, section_name, content)


def _normalize_ppt_plan(plan: dict, chapter_name: str, section_name: str, content: dict) -> dict:
    slides = plan.get('slides') if isinstance(plan.get('slides'), list) else []
    normalized = []
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_type = (slide.get('type') or 'content').strip().lower()
        title = str(slide.get('title') or '').strip()
        if slide_type == 'cover':
            normalized.append({
                'type': 'cover',
                'title': title or chapter_name,
                'subtitle': str(slide.get('subtitle') or section_name).strip(),
            })
        elif slide_type == 'summary':
            normalized.append({
                'type': 'summary',
                'title': title or '总结',
                'body': str(slide.get('body') or '').strip(),
            })
        else:
            body = str(slide.get('body') or '').strip()
            if title and body:
                normalized.append({'type': 'content', 'title': title, 'body': body})
    if not normalized:
        return _fallback_ppt_plan(chapter_name, section_name, content)
    if normalized[0]['type'] != 'cover':
        normalized.insert(0, {'type': 'cover', 'title': chapter_name, 'subtitle': section_name})
    if normalized[-1]['type'] != 'summary':
        normalized.append({
            'type': 'summary',
            'title': '总结',
            'body': f'本节《{section_name}》要点回顾，建议结合课堂练习巩固理解。',
        })
    return {
        'chapterName': chapter_name,
        'sectionName': section_name,
        'slides': normalized,
    }


def _fallback_ppt_plan(chapter_name: str, section_name: str, content: dict) -> dict:
    return {
        'chapterName': chapter_name,
        'sectionName': section_name,
        'slides': [
            {'type': 'cover', 'title': chapter_name, 'subtitle': section_name},
            {'type': 'content', 'title': '教学目标', 'body': content.get('teachingGoal', '')},
            {'type': 'content', 'title': '核心知识点', 'body': content.get('knowledgePoints', '')},
            {'type': 'content', 'title': '建议教学结构', 'body': content.get('suggestedStructure', '')},
            {'type': 'content', 'title': '其他补充信息', 'body': content.get('extraInfo', '')},
            {'type': 'summary', 'title': '总结', 'body': f'本节《{section_name}》课件已生成，可继续编辑。'},
        ],
    }


def generate_ppt_from_plan(output_path: str | Path, plan: dict) -> None:
    output_path = Path(output_path)
    ensure_dir(output_path.parent)
    prs = Presentation()
    for slide in plan.get('slides', []):
        slide_type = slide.get('type', 'content')
        if slide_type == 'cover':
            page = prs.slides.add_slide(prs.slide_layouts[0])
            page.shapes.title.text = _clip_text(slide.get('title', ''))
            if len(page.placeholders) > 1:
                page.placeholders[1].text = _clip_text(slide.get('subtitle', ''))
            continue
        page = prs.slides.add_slide(prs.slide_layouts[1])
        page.shapes.title.text = _clip_text(slide.get('title', ''))
        if len(page.placeholders) > 1:
            page.placeholders[1].text = _clip_text(slide.get('body', ''))
    prs.save(str(output_path))


def _clip_text(value: str, limit: int = 8000) -> str:
    text = str(value or '').strip()
    return text[:limit] if len(text) > limit else text


async def generate_ppt(
    output_path: str | Path,
    chapter_name: str,
    section_name: str,
    content: dict,
    output_language: int = 1,
    plan_output_path: str | Path | None = None,
) -> dict:
    plan = await plan_ppt_slides(content, chapter_name, section_name, output_language)
    if plan_output_path:
        plan_path = Path(plan_output_path)
        ensure_dir(plan_path.parent)
        plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding='utf-8')
    generate_ppt_from_plan(output_path, plan)
    return plan
